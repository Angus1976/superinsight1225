"""
PPT extractor for SuperInsight Platform.

Extracts text content from .pptx files using python-pptx,
concatenating text from all slides in order.
"""

from __future__ import annotations

import logging
import os

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

MAX_FILE_SIZE_BYTES = 100 * 1024 * 1024  # 100MB


class PPTExtractor:
    """Extracts text from .pptx files slide by slide."""

    def extract(self, file_path: str) -> str:
        """Extract all text from a .pptx file, concatenated by slide.

        Args:
            file_path: Path to the .pptx file on disk.

        Returns:
            A string with all slide text joined by newlines.

        Raises:
            RuntimeError: If python-pptx is not installed or parsing fails.
            FileNotFoundError: If file_path does not exist.
            ValueError: If file is empty, too large, or not a .pptx file.
        """
        if not PPTX_AVAILABLE:
            raise RuntimeError(
                "python-pptx is required for PPT extraction. "
                "Install with: pip install python-pptx"
            )

        self._validate_file(file_path)

        try:
            prs = Presentation(file_path)
        except Exception as exc:
            raise RuntimeError(f"Failed to open PPTX file: {exc}") from exc

        slide_texts: list[str] = []
        for slide_index, slide in enumerate(prs.slides):
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            parts.append(text)
            if parts:
                slide_texts.append("\n".join(parts))
            logger.debug("Slide %d: extracted %d text parts", slide_index + 1, len(parts))

        result = "\n\n".join(slide_texts)
        if not result.strip():
            logger.warning("No text content found in %s", file_path)
            return ""

        logger.info(
            "Extracted text from %d slides in %s (%d chars)",
            len(prs.slides), file_path, len(result),
        )
        return result

    def _validate_file(self, file_path: str) -> None:
        """Validate file existence, extension, and size."""
        if not os.path.isfile(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        if not file_path.lower().endswith(".pptx"):
            raise ValueError(
                f"Unsupported file extension: '{file_path}'. Only .pptx is supported."
            )

        file_size = os.path.getsize(file_path)
        if file_size > MAX_FILE_SIZE_BYTES:
            raise ValueError(
                f"File size ({file_size / 1024 / 1024:.1f}MB) exceeds "
                f"limit of {MAX_FILE_SIZE_BYTES / 1024 / 1024:.0f}MB"
            )
        if file_size == 0:
            raise ValueError("File is empty")
