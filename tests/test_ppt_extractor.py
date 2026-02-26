"""Unit tests for PPTExtractor."""

import os
import tempfile

import pytest
from pptx import Presentation
from pptx.util import Inches

# Import directly to avoid heavy dependency chain via __init__.py
from src.extractors.ppt import PPTExtractor


@pytest.fixture
def extractor():
    return PPTExtractor()


def _create_pptx(slides_content: list[list[str]], path: str) -> str:
    """Helper: create a .pptx file with given text per slide."""
    prs = Presentation()
    for texts in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        for i, text in enumerate(texts):
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(5), Inches(1))
            txBox.text_frame.text = text
    prs.save(path)
    return path


class TestPPTExtractorExtract:
    """Tests for PPTExtractor.extract()."""

    def test_single_slide_single_shape(self, extractor, tmp_path):
        path = str(tmp_path / "test.pptx")
        _create_pptx([["Hello World"]], path)

        result = extractor.extract(path)

        assert "Hello World" in result

    def test_multiple_slides(self, extractor, tmp_path):
        path = str(tmp_path / "test.pptx")
        _create_pptx([["Slide 1 text"], ["Slide 2 text"], ["Slide 3 text"]], path)

        result = extractor.extract(path)

        assert "Slide 1 text" in result
        assert "Slide 2 text" in result
        assert "Slide 3 text" in result

    def test_slide_order_preserved(self, extractor, tmp_path):
        path = str(tmp_path / "test.pptx")
        _create_pptx([["First"], ["Second"], ["Third"]], path)

        result = extractor.extract(path)

        assert result.index("First") < result.index("Second") < result.index("Third")

    def test_multiple_shapes_per_slide(self, extractor, tmp_path):
        path = str(tmp_path / "test.pptx")
        _create_pptx([["Shape A", "Shape B"]], path)

        result = extractor.extract(path)

        assert "Shape A" in result
        assert "Shape B" in result

    def test_empty_slides_skipped(self, extractor, tmp_path):
        """Slides with no text produce empty result without error."""
        path = str(tmp_path / "test.pptx")
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # blank, no text
        prs.save(path)

        result = extractor.extract(path)

        assert result == ""


class TestPPTExtractorValidation:
    """Tests for PPTExtractor input validation."""

    def test_file_not_found(self, extractor):
        with pytest.raises(FileNotFoundError, match="File not found"):
            extractor.extract("/nonexistent/path/file.pptx")

    def test_wrong_extension(self, extractor, tmp_path):
        path = str(tmp_path / "test.pdf")
        path_obj = tmp_path / "test.pdf"
        path_obj.write_bytes(b"fake content")

        with pytest.raises(ValueError, match="Only .pptx is supported"):
            extractor.extract(path)

    def test_empty_file(self, extractor, tmp_path):
        path = str(tmp_path / "test.pptx")
        (tmp_path / "test.pptx").write_bytes(b"")

        with pytest.raises(ValueError, match="File is empty"):
            extractor.extract(path)

    def test_corrupt_file(self, extractor, tmp_path):
        path = str(tmp_path / "test.pptx")
        (tmp_path / "test.pptx").write_bytes(b"not a real pptx file content")

        with pytest.raises(RuntimeError, match="Failed to open PPTX file"):
            extractor.extract(path)
